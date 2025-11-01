"""
Text file processor for PDF, DOCX, PPTX, MD, TXT files
"""
import os
from typing import List, Dict
import PyPDF2
from docx import Document
from pptx import Presentation
import markdown


class TextProcessor:
    """Process various text document formats"""
    
    @staticmethod
    def process_pdf(file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")
    
    @staticmethod
    def process_docx(file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip()
        except Exception as e:
            raise Exception(f"Error processing DOCX: {str(e)}")
    
    @staticmethod
    def process_pptx(file_path: str) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error processing PPTX: {str(e)}")
    
    @staticmethod
    def process_markdown(file_path: str) -> str:
        """Extract text from Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_text = file.read()
            # Convert markdown to plain text (remove formatting)
            html = markdown.markdown(md_text)
            # Simple HTML tag removal
            import re
            text = re.sub('<[^<]+?>', '', html)
            return text.strip()
        except Exception as e:
            raise Exception(f"Error processing Markdown: {str(e)}")
    
    @staticmethod
    def process_txt(file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            return text.strip()
        except Exception as e:
            raise Exception(f"Error processing TXT: {str(e)}")
    
    @staticmethod
    def process_file(file_path: str) -> Dict[str, str]:
        """
        Process any supported text file and return extracted content
        
        Args:
            file_path: Path to the file
            
        Returns:
            Dictionary with file info and extracted text
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        file_name = os.path.basename(file_path)
        
        processor_map = {
            '.pdf': TextProcessor.process_pdf,
            '.docx': TextProcessor.process_docx,
            '.pptx': TextProcessor.process_pptx,
            '.md': TextProcessor.process_markdown,
            '.txt': TextProcessor.process_txt
        }
        
        if file_ext not in processor_map:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        text = processor_map[file_ext](file_path)
        
        return {
            "file_name": file_name,
            "file_type": file_ext[1:],
            "content": text,
            "content_type": "text"
        }
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Split text into overlapping chunks
        
        Args:
            text: Text to chunk
            chunk_size: Size of each chunk
            overlap: Overlap between chunks
            
        Returns:
            List of text chunks
        """
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            start += chunk_size - overlap
            
        return chunks
